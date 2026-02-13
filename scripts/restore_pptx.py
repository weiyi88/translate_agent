#!/usr/bin/env python3
"""
PPTX 还原脚本

从 Node.js 调用，用于将翻译后的文本还原到原始 PPTX 文件中。
使用 python-pptx 库保留完整的格式信息。

使用方式:
python restore_pptx.py <original.pptx> <translations.json> <output.pptx>
"""

import json
import sys
from typing import List, Dict, Any, Union, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.dml.color import RGBColor


def get_nested_shape(slide, position: Union[int, List[int]]):
    """
    获取嵌套形状

    对应 TypeScript getNestedShape() 和 Python get_or_create_nested_shape()
    """
    if isinstance(position, int):
        return slide.shapes[position]

    current_shape = slide.shapes

    for i, index in enumerate(position):
        if i == len(position) - 1:
            # 最后一个索引
            if hasattr(current_shape, 'shapes'):
                return current_shape.shapes[index]
            else:
                return current_shape[index]

        # 中间索引 - 继续遍历
        if hasattr(current_shape, 'shapes'):
            current_shape = current_shape.shapes[index]
        else:
            current_shape = current_shape[index]

    return current_shape


def apply_font_info(text_frame, font_info: Dict[str, Any]):
    """
    应用字体信息

    对应 TypeScript applyFontInfo() 和 Python apply_font_info_safely()
    """
    try:
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            # 字体属性
            font = text_frame.paragraphs[0].runs[0].font

            if "size" in font_info and font_info["size"] is not None:
                font.size = Pt(font_info["size"] / 100)

            if "name" in font_info and font_info["name"]:
                font.name = font_info["name"]

            if "bold" in font_info and font_info["bold"] is not None:
                font.bold = font_info["bold"]

            if "italic" in font_info and font_info["italic"] is not None:
                font.italic = font_info["italic"]

            if "underline" in font_info and font_info["underline"] is not None:
                font.underline = font_info["underline"]

            if "color" in font_info and font_info["color"] is not None:
                r, g, b = font_info["color"]
                font.color.rgb = RGBColor(r, g, b)

            if "language_id" in font_info and font_info["language_id"] is not None:
                font.language_id = font_info["language_id"]

            # 段落属性
            paragraph = text_frame.paragraphs[0]

            if "alignment" in font_info and font_info["alignment"] is not None:
                paragraph.alignment = font_info["alignment"]

            if "level" in font_info and font_info["level"] is not None:
                paragraph.level = font_info["level"]

            if "line_spacing" in font_info and font_info["line_spacing"] is not None:
                paragraph.line_spacing = font_info["line_spacing"]

            if "space_before" in font_info and font_info["space_before"] is not None:
                paragraph.space_before = font_info["space_before"]

            if "space_after" in font_info and font_info["space_after"] is not None:
                paragraph.space_after = font_info["space_after"]

        # 文本框属性
        if "word_wrap" in font_info and font_info["word_wrap"] is not None:
            text_frame.word_wrap = font_info["word_wrap"]

        if "auto_size" in font_info and font_info["auto_size"] is not None:
            text_frame.auto_size = font_info["auto_size"]

        if "vertical_anchor" in font_info and font_info["vertical_anchor"] is not None:
            text_frame.vertical_anchor = font_info["vertical_anchor"]

        if "margin_left" in font_info and font_info["margin_left"] is not None:
            text_frame.margin_left = font_info["margin_left"]

        if "margin_right" in font_info and font_info["margin_right"] is not None:
            text_frame.margin_right = font_info["margin_right"]

        if "margin_top" in font_info and font_info["margin_top"] is not None:
            text_frame.margin_top = font_info["margin_top"]

        if "margin_bottom" in font_info and font_info["margin_bottom"] is not None:
            text_frame.margin_bottom = font_info["margin_bottom"]

    except Exception as e:
        print(f"Warning: Failed to apply font info: {e}", file=sys.stderr)


def restore_text_to_presentation(prs: Presentation, translations: List[Dict[str, Any]]):
    """
    还原翻译文本到演示文稿

    对应 TypeScript restoreTextToPresentation()
    """
    for element in translations:
        try:
            translated_text = element["translatedText"]
            position = element["position"]
            metadata = element["metadata"]

            slide_number, shape_index = position
            shape_type = metadata["shape_type"]
            table_info = metadata["table_info"]
            font_info = metadata["font_info"]

            # 获取幻灯片
            if slide_number >= len(prs.slides):
                print(f"Error: Slide {slide_number} not found", file=sys.stderr)
                continue

            slide = prs.slides[slide_number]

            # 获取形状
            shape = get_nested_shape(slide, shape_index)

            # 根据形状类型还原文本
            if "TABLE" in shape_type and table_info:
                # 表格单元格
                cell = shape.table.cell(table_info["row"], table_info["col"])
                cell.text = translated_text
                apply_font_info(cell.text_frame, font_info)

            elif "CHART" in shape_type:
                # 图表标题
                if hasattr(shape, 'has_title') and shape.has_title:
                    shape.chart_title.text_frame.text = translated_text
                    apply_font_info(shape.chart_title.text_frame, font_info)

            elif hasattr(shape, 'text'):
                # 普通文本形状
                shape.text = translated_text
                if hasattr(shape, 'text_frame'):
                    apply_font_info(shape.text_frame, font_info)

            else:
                print(f"Warning: Cannot handle shape type: {shape_type} at position {position}", file=sys.stderr)

        except Exception as e:
            print(f"Error restoring text at position {element.get('position')}: {e}", file=sys.stderr)


def main():
    """主函数"""
    if len(sys.argv) != 4:
        print("Usage: python restore_pptx.py <original.pptx> <translations.json> <output.pptx>", file=sys.stderr)
        sys.exit(1)

    original_path = sys.argv[1]
    translations_path = sys.argv[2]
    output_path = sys.argv[3]

    try:
        # 加载原始 PPTX
        print(f"Loading original PPTX: {original_path}", file=sys.stderr)
        prs = Presentation(original_path)

        # 加载翻译数据
        print(f"Loading translations: {translations_path}", file=sys.stderr)
        with open(translations_path, 'r', encoding='utf-8') as f:
            translations = json.load(f)

        print(f"Restoring {len(translations)} translations...", file=sys.stderr)

        # 还原翻译
        restore_text_to_presentation(prs, translations)

        # 保存输出
        print(f"Saving output: {output_path}", file=sys.stderr)
        prs.save(output_path)

        print("Success!", file=sys.stderr)
        sys.exit(0)

    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc(file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
