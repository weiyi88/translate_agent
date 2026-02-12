import logging
from typing import Any, Dict, List, Tuple, Union

from pptx import Presentation
from pptx.chart.chart import Chart
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.text.text import Font, TextFrame
from pptx.util import Pt

from app.util.config import Settings


def get_psr(file_path):
    return Presentation(file_path)


def page_num_presentation(psr: Presentation):
    return len(psr.slides)


async def save_traslate_file(psr: Presentation, out_file_name: str):
    output_path = Settings.FILE_OUTPUT_PATH+f"{out_file_name}"
    psr.save(output_path)
    return output_path


def get_or_create_nested_shape(slide, position: Union[int, Tuple]) -> Union[BaseShape, GroupShape]:
    if isinstance(position, int):
        return slide.shapes[position]

    current_shape = slide.shapes
    for i, index in enumerate(position):
        if i == len(position) - 1:
            if isinstance(current_shape, GroupShape):
                if index >= len(current_shape.shapes):
                    new_shape = current_shape.shapes.add_shape(
                        MSO_SHAPE_TYPE.TEXT_BOX, 0, 0, slide.width, slide.height)
                    return new_shape
                return current_shape.shapes[index]
            else:
                if index >= len(current_shape):
                    new_shape = current_shape.add_shape(
                        MSO_SHAPE_TYPE.TEXT_BOX, 0, 0, slide.width, slide.height)
                    return new_shape
                return current_shape[index]

        if isinstance(current_shape, GroupShape):
            if index >= len(current_shape.shapes):
                new_group = current_shape.shapes.add_group_shape()
                new_group.left = 0
                new_group.top = 0
                new_group.width = slide.width
                new_group.height = slide.height
                current_shape = new_group.shapes
            else:
                current_shape = current_shape.shapes[index]
        else:
            if index >= len(current_shape):
                new_group = current_shape.add_group_shape()
                new_group.left = 0
                new_group.top = 0
                new_group.width = slide.width
                new_group.height = slide.height
                current_shape = new_group.shapes
            else:
                current_shape = current_shape[index]
                if not isinstance(current_shape, GroupShape):
                    new_group = current_shape.add_group_shape()
                    new_group.left = 0
                    new_group.top = 0
                    new_group.width = slide.width
                    new_group.height = slide.height
                    current_shape = new_group.shapes

    return current_shape


async def restore_text_to_presentation(psr: Presentation, translated_positions: List[Tuple[str, Tuple, dict]]):
    for translated_text, position, info in translated_positions:
        try:
            slide = psr.slides[position[0]]
            shape = get_or_create_nested_shape(
                slide, position[1] if isinstance(position[1], tuple) else (position[1],))

            shape_type = info["shape_type"]
            table_info = info["table_info"]
            font_info = info["font_info"]

            if shape_type == str(MSO_SHAPE_TYPE.TABLE) and table_info:
                cell = shape.table.cell(table_info["row"], table_info["col"])
                cell.text = translated_text
                await apply_font_info_safely(cell.text_frame, font_info)
            elif shape_type == "CHART" and isinstance(shape, Chart):
                shape.chart_title.text_frame.text = translated_text
                await apply_font_info_safely(shape.chart_title.text_frame, font_info)
            elif hasattr(shape, 'text'):
                shape.text = translated_text
                await apply_font_info_safely(shape.text_frame, font_info)
            elif shape_type == "AUTO_SHAPE (1)" or shape_type == "TEXT_BOX (17)":
                if not hasattr(shape, 'text'):
                    text_box = slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height)
                    shape = text_box
                shape.text = translated_text
                await apply_font_info_safely(shape.text_frame, font_info)
            else:
                logging.warning(f"无法处理的形状类型: {shape_type}")

        except Exception as e:
            logging.error(f"还原文本时出错: {e}, 位置: {position}")


async def apply_font_info_safely(text_frame, font_info: dict):
    try:
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            font = text_frame.paragraphs[0].runs[0].font
            if "size" in font_info:
                font.size = Pt(font_info["size"] / 100)
            if "name" in font_info:
                font.name = font_info["name"]
            if "bold" in font_info:
                font.bold = font_info["bold"]
            if "italic" in font_info:
                font.italic = font_info["italic"]
            if "underline" in font_info:
                font.underline = font_info["underline"]
            if "color" in font_info and font_info["color"] is not None:
                font.color.rgb = RGBColor(*font_info["color"])
            if "language_id" in font_info:
                font.language_id = font_info["language_id"]

            paragraph = text_frame.paragraphs[0]
            if "alignment" in font_info:
                paragraph.alignment = font_info["alignment"]
            if "level" in font_info:
                paragraph.level = font_info["level"]
            if "line_spacing" in font_info:
                paragraph.line_spacing = font_info["line_spacing"]
            if "space_before" in font_info:
                paragraph.space_before = font_info["space_before"]
            if "space_after" in font_info:
                paragraph.space_after = font_info["space_after"]

            if "word_wrap" in font_info:
                text_frame.word_wrap = font_info["word_wrap"]
            if "auto_size" in font_info:
                text_frame.auto_size = font_info["auto_size"]
            if "vertical_anchor" in font_info:
                text_frame.vertical_anchor = font_info["vertical_anchor"]
            if "margin_left" in font_info:
                text_frame.margin_left = font_info["margin_left"]
            if "margin_right" in font_info:
                text_frame.margin_right = font_info["margin_right"]
            if "margin_top" in font_info:
                text_frame.margin_top = font_info["margin_top"]
            if "margin_bottom" in font_info:
                text_frame.margin_bottom = font_info["margin_bottom"]
    except Exception as e:
        logging.warning(f"应用字体信息时出错: {e}")


def extract_text_from_shape(shape, slide_number: int, shape_index: int) -> List[Tuple[str, Tuple[int, int], Dict[str, Any]]]:
    result = []
    font_info = {}

    if hasattr(shape, 'text') and shape.text.strip():
        font_info = get_font_info_safely(shape)
        result.append([shape.text.strip(), (slide_number, shape_index),
                       {"shape_type": str(shape.shape_type), "table_info": None, "font_info": font_info}])

    match shape.shape_type:
        case MSO_SHAPE_TYPE.TABLE:
            for row_index, row in enumerate(shape.table.rows):
                for col_index, cell in enumerate(row.cells):
                    if cell.text.strip():
                        font_info = get_font_info_safely(cell)
                        result.append([cell.text.strip(), (slide_number, shape_index),
                                       {"shape_type": str(shape.shape_type), "table_info": {"row": row_index, "col": col_index}, "font_info": font_info}])
        case _:
            if isinstance(shape, Chart):
                if shape.has_title and shape.chart_title.text_frame.text.strip():
                    font_info = get_font_info_safely(shape.chart_title)
                    result.append([shape.chart_title.text_frame.text.strip(), (slide_number, shape_index),
                                   {"shape_type": "CHART", "table_info": None, "font_info": font_info}])

    return result


def extract_text_from_group(group, slide_number: int, group_index: int) -> List[Tuple[str, Tuple[int, int], Dict[str, Any]]]:
    result = []
    for sub_shape_index, sub_shape in enumerate(group.shapes):
        match sub_shape.shape_type:
            case MSO_SHAPE_TYPE.GROUP:
                result.extend(extract_text_from_group(
                    sub_shape, slide_number, (*group_index, sub_shape_index)))
            case _:
                sub_result = extract_text_from_shape(
                    sub_shape, slide_number, (*group_index, sub_shape_index))
                result.extend(sub_result)
    return result


async def extract_text_from_presentation(psr: Presentation):
    for slide_number, slide in enumerate(psr.slides):
        logging.info(f"--------------------目前进度页数: {slide_number}")
        tem_arr = []
        for shape_index, shape in enumerate(slide.shapes):
            match shape.shape_type:
                case MSO_SHAPE_TYPE.GROUP:
                    tem_arr.extend(extract_text_from_group(
                        shape, slide_number, (shape_index,)))
                case _:
                    tem_arr.extend(extract_text_from_shape(
                        shape, slide_number, shape_index))
        yield tem_arr


def get_font_info_safely(shape):
    try:
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            return get_font_info(shape.text_frame, shape.text_frame.paragraphs[0].runs[0].font)
    except Exception as e:
        logging.warning(f"获取字体信息时出错: {e}")
    return {}


def get_font_info(text_frame: TextFrame, font: Font):
    font_info = {}
    try:
        if font.size:
            font_info["size"] = font.size.pt * 100
        if font.name:
            font_info["name"] = font.name
        if font.bold is not None:
            font_info["bold"] = font.bold
        if font.italic is not None:
            font_info["italic"] = font.italic
        if font.underline is not None:
            font_info["underline"] = font.underline
        if font.color and font.color.rgb:
            font_info["color"] = font.color.rgb
        if font.language_id is not None:
            font_info["language_id"] = font.language_id

        paragraph = text_frame.paragraphs[0]
        font_info["alignment"] = paragraph.alignment
        font_info["level"] = paragraph.level
        if paragraph.line_spacing:
            font_info["line_spacing"] = paragraph.line_spacing
        font_info["space_before"] = paragraph.space_before
        font_info["space_after"] = paragraph.space_after

        font_info["word_wrap"] = text_frame.word_wrap
        font_info["auto_size"] = text_frame.auto_size
        font_info["vertical_anchor"] = text_frame.vertical_anchor
        font_info["margin_left"] = text_frame.margin_left
        font_info["margin_right"] = text_frame.margin_right
        font_info["margin_top"] = text_frame.margin_top
        font_info["margin_bottom"] = text_frame.margin_bottom

    except Exception as e:
        logging.warning(f"获取字体信息时出错: {e}")
    return font_info


async def chunk_content_by_lenth(all_text: list, target_length=Settings.CHUNK_LENGTH) -> list:
    chunked_text = []
    text_length = 0
    temp_arr = []
    for content in all_text:
        if text_length > target_length:
            chunked_text.append(temp_arr)
            temp_arr = []
            text_length = 0

        text_length += len(content[0])
        temp_arr.append(content)

    chunked_text.append(temp_arr)
    return chunked_text
